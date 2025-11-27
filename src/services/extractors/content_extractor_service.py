import os
import logging
from typing import Optional

from docx import Document
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_pdf_text

logger = logging.getLogger(__name__)


class ContentExtractorService:
    """
    Extracts text from multiple file types:
    - PDF
    - DOCX
    - PPTX
    - TXT
    """

    SUPPORTED_EXTENSIONS = ["pdf", "docx", "pptx", "txt"]

    def extract_text(self, file_path: str) -> Optional[str]:
        """
        Detect file type and extract text accordingly.
        Returns extracted text or None.
        """

        if not file_path or not os.path.exists(file_path):
            logger.error(f"❌ File does not exist: {file_path}")
            return None

        ext = file_path.split(".")[-1].lower()

        try:
            if ext == "pdf":
                return self._extract_pdf(file_path)

            elif ext == "docx":
                return self._extract_docx(file_path)

            elif ext == "pptx":
                return self._extract_pptx(file_path)

            elif ext == "txt":
                return self._extract_txt(file_path)

            else:
                logger.warning(f"⚠ Unsupported file format for extraction: {ext}")
                return None

        except Exception as e:
            logger.error(f"❌ Error extracting text from {file_path}: {e}")
            return None


    # ---------------------------------------------------------
    # PDF Extraction
    # ---------------------------------------------------------
    def _extract_pdf(self, file_path: str) -> str:
        try:
            text = extract_pdf_text(file_path)
            return text.strip()
        except Exception as e:
            logger.error(f"❌ PDF extraction failed: {e}")
            return ""


    # ---------------------------------------------------------
    # DOCX Extraction
    # ---------------------------------------------------------
    def _extract_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            lines = []

            for p in doc.paragraphs:
                if p.text.strip():
                    lines.append(p.text.strip())

            return "\n".join(lines)

        except Exception as e:
            logger.error(f"❌ DOCX extraction failed: {e}")
            return ""


    # ---------------------------------------------------------
    # PPTX Extraction
    # ---------------------------------------------------------
    def _extract_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            output = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            output.append(text)

            return "\n".join(output)

        except Exception as e:
            logger.error(f"❌ PPTX extraction failed: {e}")
            return ""


    # ---------------------------------------------------------
    # TXT Extraction
    # ---------------------------------------------------------
    def _extract_txt(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()
        except Exception as e:
            logger.error(f"❌ TXT extraction failed: {e}")
            return ""
