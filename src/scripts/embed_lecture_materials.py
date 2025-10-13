

# import argparse, pathlib, logging, sys
# from docx import Document
# import pdfplumber
# from dotenv import load_dotenv

# from src.utils.token_chunker import chunk_text
# from src.models.lecture_chunk import LectureChunk
# from src.services.embedding.openai_embedder import OpenAIEmbedder
# from src.services.embedding.gemini_embedder import GeminiEmbedder
# from src.services.database_services.lecture_material_embedding_db import LectureMaterialEmbeddingDB

# load_dotenv()
# logging.basicConfig(level=logging.INFO)
# log = logging.getLogger(__name__)

# # ────────────────────────────────────────────────
# def read_text(path: pathlib.Path) -> str:
#     if path.suffix.lower() == ".docx":
#         doc = Document(path)
#         return "\n".join(p.text for p in doc.paragraphs)
#     if path.suffix.lower() == ".pdf":
#         with pdfplumber.open(path) as pdf:
#             return "\n".join(p.extract_text() or "" for p in pdf.pages)
#     raise ValueError(f"Unsupported file type: {path.name}")

# # ────────────────────────────────────────────────
# def main() -> None:
#     ap = argparse.ArgumentParser(description="Embed lecture materials into vector DB")
#     ap.add_argument("--provider", required=True, choices=["OpenAI", "GoogleGemini"], help="Embedding provider")
#     ap.add_argument("--model", required=True, help="LLM model name (unused here, for logging)")
#     ap.add_argument("--embedder", default="text-embedding-3-small", help="Embedding model ID (e.g., text-embedding-3-small or embedding-001)")
#     ap.add_argument("--root", default="data/Lecture_Material", help="Root directory containing module folders")
#     ap.add_argument("--module", help="Target module code (e.g., EE6250)")
#     ap.add_argument("--max_tokens", type=int, default=1000)
#     ap.add_argument("--overlap", type=int, default=200)
#     args = ap.parse_args()

#     # Prompt for module if not provided
#     if not args.module:
#         print("📚 Available modules under", args.root)
#         for p in pathlib.Path(args.root).iterdir():
#             if p.is_dir():
#                 print(" •", p.name)
#         args.module = input("\nEnter module code to embed: ").strip()

#     module_dir = pathlib.Path(args.root) / args.module
#     if not module_dir.exists():
#         log.error("❌ Folder does not exist: %s", module_dir)
#         sys.exit(1)

#     # Choose embedder
#     embedder = (
#         OpenAIEmbedder(args.embedder)
#         if args.provider == "OpenAI"
#         else GeminiEmbedder(model_name=args.embedder)
#     )

#     # Create vector DB instance
#     vec_db = LectureMaterialEmbeddingDB(embedder)

#     # Scan files
#     files = list(module_dir.rglob("*"))
#     valid_files = [f for f in files if f.suffix.lower() in {".pdf", ".docx"}]
#     if not valid_files:
#         log.warning("⚠️ No supported files found in %s", module_dir)
#         sys.exit(0)

#     for f in valid_files:
#         file_name = f.name
#         if vec_db.document_exists(args.module, file_name):
#             log.info("⚠️ Skipping already embedded document: %s", file_name)
#             continue

#         log.info("📄 Processing: %s", f.relative_to(module_dir))
#         try:
#             raw_text = read_text(f)
#             chunks = chunk_text(raw_text, max_tokens=args.max_tokens, overlap=args.overlap)

#             lecture_chunks = [
#                 LectureChunk(module_code=args.module, source_file=file_name, chunk_id=i, text=txt)
#                 for i, txt in enumerate(chunks) if txt.strip()
#             ]

#             log.info("   → Saving %d chunks", len(lecture_chunks))
#             vec_db.save_chunks(lecture_chunks)

#         except Exception as e:
#             log.error("   ❌ Error processing %s: %s", f.name, str(e))

#     vec_db.close()
#     print(f"✅ Finished embedding for module: {args.module}")

# # ────────────────────────────────────────────────
# if __name__ == "__main__":
#     main()

import argparse, pathlib, logging, sys
from docx import Document
import pdfplumber
from dotenv import load_dotenv
import psycopg2
from psycopg2.extras import RealDictCursor

from src.utils.token_chunker import chunk_text
from src.models.lecture_chunk import LectureChunk
from src.services.embedding.openai_embedder import OpenAIEmbedder
from src.services.embedding.gemini_embedder import GeminiEmbedder
from src.services.database_services.lecture_material_embedding_db import (
    LectureMaterialEmbeddingDB,
)

load_dotenv()
logging.basicConfig(level=logging.INFO)
log = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────────────
def get_database_connection():
    """Get database connection using environment variables."""
    try:
        conn = psycopg2.connect(
            host=os.getenv('POSTGRES_HOST'),
            port=os.getenv('POSTGRES_PORT'),
            database=os.getenv('POSTGRES_DB'),
            user=os.getenv('POSTGRES_USER'),
            password=os.getenv('POSTGRES_PASSWORD')
        )
        return conn
    except Exception as e:
        log.error(f"Failed to connect to database: {e}")
        sys.exit(1)


# ──────────────────────────────────────────────────────────────────────
def get_assessment_related_lecture_materials(assessment_id):
    """Retrieve lecture materials related to specific assessment through module and lessons."""
    conn = get_database_connection()
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT DISTINCT lm.material_id, lm.lesson_id, lm.file_name, lm.file_url, 
                       lm.uploaded_on, lm.description,
                       l.lesson_id, l.title as lesson_title,
                       m.module_id, m.module_code, m.module_name
                FROM "Assessment" a
                JOIN "Module" m ON a.module_id = m.module_id
                JOIN "Lesson" l ON m.module_id = l.module_id
                JOIN "LectureMaterial" lm ON l.lesson_id = lm.lesson_id
                WHERE a.assessment_id = %s
                ORDER BY lm.uploaded_on ASC
            """, (assessment_id,))
            
            materials = cur.fetchall()
            log.info(f"Found {len(materials)} lecture materials for assessment {assessment_id}")
            return materials
            
    except Exception as e:
        log.error(f"Database error: {e}")
        return []
    finally:
        conn.close()


def get_all_lecture_materials_from_db(module_filter=None):
    """Retrieve lecture material file paths from database."""
    conn = get_database_connection()
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            if module_filter:
                cur.execute("""
                    SELECT lm.material_id, lm.lesson_id, lm.file_name, lm.file_url, 
                           lm.uploaded_on, lm.description,
                           l.module_id
                    FROM "LectureMaterial" lm
                    JOIN "Lesson" l ON lm.lesson_id = l.lesson_id
                    WHERE l.module_id = %s
                    ORDER BY l.module_id, lm.uploaded_on ASC
                """, (module_filter,))
            else:
                cur.execute("""
                    SELECT lm.material_id, lm.lesson_id, lm.file_name, lm.file_url, 
                           lm.uploaded_on, lm.description,
                           l.module_id
                    FROM "LectureMaterial" lm
                    JOIN "Lesson" l ON lm.lesson_id = l.lesson_id
                    ORDER BY l.module_id, lm.uploaded_on ASC
                """)
            
            materials = cur.fetchall()
            log.info(f"Found {len(materials)} lecture materials in database" + 
                    (f" for module {module_filter}" if module_filter else ""))
            return materials
            
    except Exception as e:
        log.error(f"Database error: {e}")
        return []
    finally:
        conn.close()


# ──────────────────────────────────────────────────────────────────────
def get_available_modules():
    """Get list of available modules from database."""
    conn = get_database_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT DISTINCT l.module_id 
                FROM "Lesson" l
                JOIN "LectureMaterial" lm ON l.lesson_id = lm.lesson_id
                ORDER BY l.module_id
            """)
            modules = [row[0] for row in cur.fetchall()]
            return modules
    except Exception as e:
        log.error(f"Error fetching modules: {e}")
        return []
    finally:
        conn.close()


# ──────────────────────────────────────────────────────────────────────
def get_project_root() -> pathlib.Path:
    """Get the project root directory (current working directory)."""
    return pathlib.Path.cwd()


# ──────────────────────────────────────────────────────────────────────
def read_text(file_path: str) -> str:
    """Read text from PDF, DOCX, or PowerPoint file."""
    # The file_path from database is relative to the parent of project root
    project_root = get_project_root()
    full_path = project_root.parent / file_path
    
    log.info(f"Looking for file at: {full_path}")
    
    if not full_path.exists():
        # Let's also check if it might be in the project directory itself
        alternative_path = project_root / file_path
        log.info(f"Alternative path: {alternative_path}")
        
        if alternative_path.exists():
            full_path = alternative_path
        else:
            # List what's actually in the expected directory
            expected_dir = full_path.parent
            if expected_dir.exists():
                files_in_dir = list(expected_dir.iterdir())
                log.error(f"Expected directory exists but file not found. Contents: {[f.name for f in files_in_dir]}")
            else:
                log.error(f"Expected directory does not exist: {expected_dir}")
            raise FileNotFoundError(f"File not found at: {full_path} or {alternative_path}")
    
    log.info(f"Successfully found file at: {full_path}")
    
    if full_path.suffix.lower() == ".docx":
        doc = Document(full_path)
        return "\n".join(p.text for p in doc.paragraphs)
    elif full_path.suffix.lower() == ".pdf":
        with pdfplumber.open(full_path) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    elif full_path.suffix.lower() in [".ppt", ".pptx"]:
        # For PowerPoint files, we need python-pptx library
        try:
            from pptx import Presentation
            prs = Presentation(full_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return "\n".join(text_content)
        except ImportError:
            raise ImportError("python-pptx library is required for PowerPoint files. Install with: pip install python-pptx")
    else:
        raise ValueError(f"Unsupported file type: {full_path.name}")


# ──────────────────────────────────────────────────────────────────────
def get_document_identifier(material):
    """Get a unique identifier for the document to check for duplicates."""
    # Use file_name if available, otherwise use description or filename from path
    if material.get('file_name'):
        return material['file_name']
    elif material.get('description'):
        return material['description']
    else:
        return pathlib.Path(material['file_url']).name


# ──────────────────────────────────────────────────────────────────────
def main() -> None:
    ap = argparse.ArgumentParser(description="Embed lecture materials into vector DB")
    ap.add_argument("--provider", required=True, choices=["OpenAI", "GoogleGemini", "DeepSeek"], help="Embedding provider")
    ap.add_argument("--model", required=True, help="LLM model name (unused here, for logging)")
    ap.add_argument("--embedder", default="text-embedding-3-small", help="Embedding model ID (e.g., text-embedding-3-small or embedding-001)")
    ap.add_argument("--root", default="data/Lecture_Material", help="Root directory containing module folders")
    ap.add_argument("--module", help="Target module code (e.g., EE6250)")
    ap.add_argument("--max_tokens", type=int, default=1000)
    ap.add_argument("--overlap", type=int, default=200)
    args = ap.parse_args()

    # Get lecture materials based on filters
    materials = []
    
    if args.assessment_id:
        # Filter lecture materials by assessment
        log.info(f"Processing lecture materials for assessment: {args.assessment_id}")
        materials = get_assessment_related_lecture_materials(args.assessment_id)
        
        if not materials:
            log.error(f"❌ No lecture materials found for assessment {args.assessment_id}")
            sys.exit(1)
            
    elif args.module:
        # Filter by module
        available_modules = get_available_modules()
        if args.module not in available_modules:
            log.error(f"❌ Module '{args.module}' not found in database")
            sys.exit(1)
        log.info(f"Processing specific module: {args.module}")
        materials = get_all_lecture_materials_from_db(args.module)
        
    else:
        # Process all modules
        log.info("Processing all modules from database")
        materials = get_all_lecture_materials_from_db()

    if not materials:
        log.error("❌ No lecture materials found in database. Please upload some files first.")
        sys.exit(1)

    # Group materials by module for better logging
    modules = {}
    for material in materials:
        module_key = material.get('module_code') or material.get('module_id', 'unknown')
        if module_key not in modules:
            modules[module_key] = []
        modules[module_key].append(material)

    log.info(f"Found {len(materials)} files across {len(modules)} modules: {', '.join(modules.keys())}")

    # Choose embedder
    provider_override = None
    if args.provider == "OpenAI":
        embedder = OpenAIEmbedder(args.embedder)
    elif args.provider == "GoogleGemini":
        embedder = GeminiEmbedder(model_name=args.embedder)
    elif args.provider == "DeepSeek":
        print("⚠️ DeepSeek selected: using OpenAI embeddings but saving in DeepSeek table")
        embedder = OpenAIEmbedder("text-embedding-3-small")
        provider_override = "deepseek"
    else:
        raise ValueError(f"Unsupported provider: {args.provider}")

    # Create vector DB instance
    vec_db = LectureMaterialEmbeddingDB(embedder, provider_override=provider_override)

    # Process each material
    total_chunks = 0
    processed_files = 0
    skipped_files = 0
    
    for material in materials:
        file_path = material['file_url']
        file_identifier = get_document_identifier(material)
        module_key = material.get('module_code') or material.get('module_id', 'unknown')
        
        # Check if document already exists to avoid duplicates
        # Pass assessment_id to check for existing embeddings
        if vec_db.document_exists(module_key, file_identifier, args.assessment_id):
            log.info(f"⚠️ [{module_key}] Skipping already embedded document: {file_identifier} (Assessment: {args.assessment_id})")
            skipped_files += 1
            continue
        
        try:
            log.info(f"📄 [{module_key}] Processing: {file_identifier}")
            raw_text = read_text(file_path)
            
            if not raw_text.strip():
                log.warning(f"⚠️ No text extracted from {file_identifier}")
                continue
            
            # Create chunks
            chunks = chunk_text(raw_text, max_tokens=args.max_tokens, overlap=args.overlap)
            
            # Create LectureChunk objects
            lecture_chunks = [
                LectureChunk(
                    assessment_id=args.assessment_id,
                    module_code=module_key, 
                    source_file=file_identifier, 
                    chunk_id=idx, 
                    text=txt
                )
                for idx, txt in enumerate(chunks) if txt.strip()
            ]
            
            if lecture_chunks:
                log.info(f"   → Saving {len(lecture_chunks)} chunks")
                # Pass assessment_id when saving chunks
                vec_db.save_chunks(lecture_chunks, args.assessment_id)
                total_chunks += len(lecture_chunks)
                processed_files += 1
            else:
                log.warning(f"⚠️ No valid chunks created from {file_identifier}")
                
        except Exception as e:
            log.error(f"❌ Error processing {file_identifier}: {e}")
            continue

    vec_db.close()
    
    # Final summary
    print(f"\n✅ Embedding Summary:")
    print(f"   • Processed files: {processed_files}")
    print(f"   • Skipped files: {skipped_files}")
    print(f"   • Total chunks created: {total_chunks}")
    print(f"   • Modules: {', '.join(modules.keys())}")
    if args.assessment_id:
        print(f"   • Assessment ID: {args.assessment_id}")
    
    if processed_files > 0:
        print(f"🎉 Successfully embedded {processed_files} files with {total_chunks} chunks!")
    else:
        print("⚠️ No new files were processed")


if __name__ == "__main__":
    main()